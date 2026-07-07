import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import os

def pdf_to_pptx(pdf_path, pptx_path):
    print(f"Converting {pdf_path} to {pptx_path}...")
    doc = fitz.open(pdf_path)
    prs = Presentation()

    # Define standard slide size (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        
        # Determine the scale based on slide size (to get good resolution)
        zoom = 2.0  # increase for better resolution
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat)
        
        image_path = f"temp_page_{page_num}.png"
        pix.save(image_path)
        
        blank_slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Calculate scaling to fit slide
        width = prs.slide_width
        height = prs.slide_height
        
        slide.shapes.add_picture(image_path, 0, 0, width=width, height=height)
        
        os.remove(image_path)
        print(f"Processed page {page_num + 1}/{len(doc)}")

    prs.save(pptx_path)
    print("Conversion complete!")

import sys

if __name__ == "__main__":
    if len(sys.argv) > 1:
        pdf_file = sys.argv[1]
        pptx_file = pdf_file.rsplit(".", 1)[0] + ".pptx"
    else:
        pdf_file = "/Users/nicksng/code/random/Final 3mn - 5mn CheckinMe pitch 2026.pdf"
        pptx_file = "/Users/nicksng/code/random/Final 3mn - 5mn CheckinMe pitch 2026.pptx"
    pdf_to_pptx(pdf_file, pptx_file)
