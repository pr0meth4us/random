import os
from pptx import Presentation
from pptx.util import Inches

def fix_presentation(pptx_path):
    print(f"Fixing {pptx_path}...")
    prs = Presentation(pptx_path)
    
    # Create a new presentation to copy into
    new_prs = Presentation()
    new_prs.slide_width = prs.slide_width
    new_prs.slide_height = prs.slide_height
    
    for i, slide in enumerate(prs.slides):
        image_part = None
        # Find the background image in the slide relations
        for rel in slide.part.rels.values():
            if "image" in rel.target_ref:
                image_part = rel.target_part
                break
                
        blank_slide_layout = new_prs.slide_layouts[6]
        new_slide = new_prs.slides.add_slide(blank_slide_layout)
        
        if image_part:
            image_bytes = image_part.blob
            image_path = f"temp_bg_{i}.png"
            with open(image_path, "wb") as f:
                f.write(image_bytes)
            
            new_slide.shapes.add_picture(
                image_path, 0, 0, 
                width=new_prs.slide_width, 
                height=new_prs.slide_height
            )
            os.remove(image_path)
            
    fixed_path = pptx_path.replace(".pptx", "_fixed.pptx")
    new_prs.save(fixed_path)
    print(f"Fixed presentation saved to {fixed_path}!")

if __name__ == "__main__":
    fix_presentation("/Users/nicksng/code/random/7. EZEPOS.pptx")
