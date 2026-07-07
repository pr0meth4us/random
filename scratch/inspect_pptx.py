import sys
from pptx import Presentation

def inspect_pptx(pptx_path):
    print(f"Opening {pptx_path}...")
    prs = Presentation(pptx_path)
    
    print(f"Dimensions: {prs.slide_width.inches:.2f} x {prs.slide_height.inches:.2f} inches")
    print(f"Total slides: {len(prs.slides)}\n")
    
    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i+1} ---")
        shapes = slide.shapes
        print(f"Number of shapes: {len(shapes)}")
        
        # Check relations for images
        images = []
        for rel in slide.part.rels.values():
            if "image" in rel.target_ref:
                images.append(rel.target_part)
        print(f"Number of related images: {len(images)}")
        if images:
            for idx, img in enumerate(images):
                print(f"  Image {idx+1}: {img.content_type}, size={len(img.blob)} bytes")
                
        # Extract text from shapes
        text_runs = []
        for shape in shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
            # Table text
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text_runs.append(cell.text)
                        
        if text_runs:
            print("Extracted Text:")
            for run_text in text_runs:
                if run_text.strip():
                    print(f"  - {run_text.strip()}")
        else:
            print("No text shapes found on this slide.")
        
        # Notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                print(f"Notes:\n  {notes}")
        print()

if __name__ == "__main__":
    pptx_path = "/Users/nicksng/code/random/EGD_Slide Presentaton_DA5.pptx"
    inspect_pptx(pptx_path)
