import collections
import collections.abc
from pptx import Presentation
import hashlib

prs = Presentation('EGD_Slide Presentaton_DA5.pptx')

s60 = prs.slides[59] # 0-indexed
s72 = prs.slides[71] # 0-indexed

def get_text_and_image_hashes(slide):
    text_content = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_content.append(run.text)
        elif shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    text_content.append(cell.text)
    
    text_str = "".join(text_content)
    text_hash = hashlib.md5(text_str.encode('utf-8')).hexdigest()
    
    image_hashes = []
    for rel in slide.part.rels.values():
        if "image" in rel.target_ref:
            img_hash = hashlib.md5(rel.target_part.blob).hexdigest()
            image_hashes.append(img_hash)
            
    return text_hash, text_str, image_hashes

h60_txt, s60_txt, h60_imgs = get_text_and_image_hashes(s60)
h72_txt, s72_txt, h72_imgs = get_text_and_image_hashes(s72)

print("Slide 60:")
print(f"  Text MD5: {h60_txt}")
print(f"  Images MD5s: {h60_imgs}")
print("Slide 72:")
print(f"  Text MD5: {h72_txt}")
print(f"  Images MD5s: {h72_imgs}")

if s60_txt != s72_txt:
    print("Texts differ!")
    print(f"S60 text len: {len(s60_txt)}")
    print(f"S72 text len: {len(s72_txt)}")
