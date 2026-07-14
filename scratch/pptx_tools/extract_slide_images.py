import os
from pptx import Presentation

def extract_images(pptx_path, output_dir):
    os.makedirs(output_dir, exist_ok=True)
    prs = Presentation(pptx_path)
    
    print(f"Opening {pptx_path}...")
    
    # We want to extract images for slides that have 1 image and no text, or slides of interest
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        
        # Extract text first to check
        text_content = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_content.append(run.text)
        has_text = len("".join(text_content).strip()) > 0
        
        # Extract images
        image_parts = []
        for rel in slide.part.rels.values():
            if "image" in rel.target_ref:
                image_parts.append(rel.target_part)
                
        if not has_text and len(image_parts) == 1:
            img = image_parts[0]
            # Determine extension
            ext = "png" if "png" in img.content_type else "jpg"
            if "jpeg" in img.content_type:
                ext = "jpg"
            elif "gif" in img.content_type:
                ext = "gif"
                
            out_path = os.path.join(output_dir, f"slide_{slide_num:03d}.{ext}")
            with open(out_path, "wb") as f:
                f.write(img.blob)
            print(f"Slide {slide_num:03d}: Extracted image to {out_path} ({len(img.blob)} bytes)")

if __name__ == "__main__":
    pptx_path = "/Users/nicksng/code/random/EGD_Slide Presentaton_DA5.pptx"
    output_dir = "/Users/nicksng/code/random/scratch/slide_images"
    extract_images(pptx_path, output_dir)
