import os
import sys
from pptx import Presentation
from google import genai

# Setup API Key
os.environ["GEMINI_API_KEY"] = os.getenv("GEMINI_API_KEY", "")
client = genai.Client()

def extract_slide_text(slide):
    text_runs = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
        elif shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    text_runs.append(cell.text)
    return "\n".join([line.strip() for line in text_runs if line.strip()])

def review_text_slides(pptx_path, output_path):
    prs = Presentation(pptx_path)
    slides_text = []
    
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        text = extract_slide_text(slide)
        if text:
            slides_text.append((slide_num, text))
            
    print(f"Found {len(slides_text)} slides with text. Processing...")
    
    # We will batch slides in groups of 10
    batch_size = 10
    with open(output_path, "w", encoding="utf-8") as out_f:
        out_f.write("=== EDITABLE TEXT SLIDES REVIEW ===\n\n")
        
        for idx in range(0, len(slides_text), batch_size):
            batch = slides_text[idx:idx+batch_size]
            prompt = "Review the following PowerPoint slides for spelling, grammatical, and typographical errors. The slides are written in Khmer and English.\n"
            prompt += "Look for:\n"
            prompt += "1. Khmer spelling errors or typos (e.g. wrong vowels, missing subscripts, incorrect word spacing/breaks).\n"
            prompt += "2. English spelling errors or typos.\n"
            prompt += "3. Awkward phrasing, syntax issues, or broken sentences.\n\n"
            prompt += "Format your response as a list of findings, specifying for each:\n"
            prompt += "- Slide number\n"
            prompt += "- Original (Incorrect) text\n"
            prompt += "- Corrected text\n"
            prompt += "- Reason for correction (in English)\n"
            prompt += "If a slide has no errors, do not list it. If a batch has no errors at all, reply 'No errors found in this batch.'\n\n"
            
            for s_num, s_text in batch:
                prompt += f"--- Slide {s_num} ---\n{s_text}\n\n"
                
            print(f"Sending batch {idx // batch_size + 1} (Slides {batch[0][0]} to {batch[-1][0]})...")
            
            try:
                response = client.models.generate_content(
                    model='gemini-2.5-flash',
                    contents=prompt
                )
                out_f.write(f"--- Batch {idx // batch_size + 1} (Slides {batch[0][0]} to {batch[-1][0]}) ---\n")
                out_f.write(response.text)
                out_f.write("\n\n" + "="*50 + "\n\n")
            except Exception as e:
                print(f"Error in batch: {e}")
                out_f.write(f"--- Batch {idx // batch_size + 1} ERROR: {e} ---\n\n")
                
    print(f"Review completed. Results saved to {output_path}")

if __name__ == "__main__":
    pptx_path = "/Users/nicksng/code/random/EGD_Slide Presentaton_DA5.pptx"
    output_path = "/Users/nicksng/code/random/scratch/text_review_results.txt"
    review_text_slides(pptx_path, output_path)
