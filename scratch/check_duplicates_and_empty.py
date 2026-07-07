from pptx import Presentation
import hashlib

def get_slide_hash(slide):
    # Hash the text and image bytes of the slide to find duplicates
    hasher = hashlib.md5()
    
    # 1. Add all texts
    text_content = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_content.append(run.text)
        elif shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    text_content.append(cell.text)
    hasher.update("".join(text_content).encode('utf-8'))
    
    # 2. Add all image sizes/blobs if they exist
    images = []
    for rel in slide.part.rels.values():
        if "image" in rel.target_ref:
            # Hash the actual image blob
            hasher.update(rel.target_part.blob)
            
    return hasher.hexdigest(), "".join(text_content)

def check_presentation(pptx_path):
    prs = Presentation(pptx_path)
    
    slide_hashes = {}
    duplicates = []
    
    image_only_slides = []
    
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        h, txt = get_slide_hash(slide)
        
        # Check if duplicate
        if h in slide_hashes:
            duplicates.append((slide_num, slide_hashes[h]))
        else:
            slide_hashes[h] = slide_num
            
        # Check if image-only (no text, but has images)
        shapes = slide.shapes
        has_text = len(txt.strip()) > 0
        
        has_image = False
        images_count = 0
        for rel in slide.part.rels.values():
            if "image" in rel.target_ref:
                has_image = True
                images_count += 1
                
        if not has_text and has_image:
            image_only_slides.append((slide_num, images_count))
            
    print(f"Total Slides Checked: {len(prs.slides)}")
    print("\n--- Duplicate Slides Found ---")
    if duplicates:
        for dup, orig in duplicates:
            print(f"Slide {dup} is a DUPLICATE of Slide {orig}")
    else:
        print("No exact duplicates found.")
        
    print("\n--- Image-Only Slides (No editable text, contains images) ---")
    if image_only_slides:
        for s_num, img_c in image_only_slides:
            print(f"Slide {s_num} has no text shapes, but has {img_c} image(s)")
    else:
        print("No image-only slides found.")

if __name__ == "__main__":
    pptx_path = "/Users/nicksng/code/random/EGD_Slide Presentaton_DA5.pptx"
    check_presentation(pptx_path)
