import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("python-pptx is not installed. Please run: pip install python-pptx")
    sys.exit(1)

def extract_text_from_pptx(pptx_path: str):
    input_path = Path(pptx_path).resolve()
    
    if not input_path.exists():
        print(f"Error: {input_path} does not exist.")
        sys.exit(1)

    print(f"Extracting text locally from {input_path.name}...")
    
    try:
        prs = Presentation(input_path)
    except Exception as e:
        print(f"Failed to open presentation: {e}")
        sys.exit(1)

    extracted_text = []
    
    for i, slide in enumerate(prs.slides):
        extracted_text.append(f"--- Slide {i + 1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                extracted_text.append(shape.text.strip())
        extracted_text.append("\n")

    out_txt = input_path.with_name(f"{input_path.stem}_transcription.txt")
    with open(out_txt, "w", encoding="utf-8") as f:
        f.write("\n".join(extracted_text))
        
    print(f"✅ Text extraction complete! Saved to: {out_txt}")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_pptx_text.py <file.pptx>")
        sys.exit(1)
    extract_text_from_pptx(sys.argv[1])
