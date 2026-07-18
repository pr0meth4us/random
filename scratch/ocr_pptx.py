import os
import sys
import time
import io
from pathlib import Path
from pptx import Presentation

try:
    from google.cloud import vision
except ImportError:
    print("google-cloud-vision is not installed. Please run: uv pip install --system google-cloud-vision")
    sys.exit(1)

def detect_text(content: bytes) -> str:
    """Detects document text in an image using Google Cloud Vision API."""
    client = vision.ImageAnnotatorClient()
    image = vision.Image(content=content)
    
    # document_text_detection is optimized for dense text and documents
    response = client.document_text_detection(image=image)
    
    if response.error.message:
        raise Exception(f"Vision API Error: {response.error.message}")
        
    if response.full_text_annotation:
        return response.full_text_annotation.text
    return ""

def process_pptx(pptx_path: str):
    input_path = Path(pptx_path).resolve()
    
    if not input_path.exists():
        print(f"Error: {input_path} does not exist.")
        sys.exit(1)

    print(f"Opening presentation: {input_path.name}")
    try:
        prs = Presentation(input_path)
    except Exception as e:
        print(f"Failed to open presentation: {e}")
        sys.exit(1)

    out_txt = input_path.with_name(f"{input_path.stem}_vision_hybrid_transcription.txt")
    
    with open(out_txt, "w", encoding="utf-8") as f:
        f.write(f"--- Vision OCR Hybrid Transcription for {input_path.name} ---\n\n")

    print(f"Total slides to process: {len(prs.slides)}")

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        print(f"\n--- Processing Slide {slide_num}/{len(prs.slides)} ---")
        
        # Only OCR specific slides
        needs_ocr = (19 <= slide_num <= 45) or (67 <= slide_num <= 70)
        
        if not needs_ocr:
            extracted_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    extracted_text.append(shape.text.strip())
            raw_text = "\n".join(extracted_text)
            
            with open(out_txt, "a", encoding="utf-8") as f:
                f.write(f"--- Slide {slide_num} ---\n")
                f.write(f"{raw_text}\n\n")
            print(f"Slide {slide_num}: ✅ Raw text extracted.")
            continue
            
        # OCR Path
        largest_img_blob = None
        max_size = 0
        
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                size = len(shape.image.blob)
                if size > max_size:
                    max_size = size
                    largest_img_blob = shape.image.blob
        
        if largest_img_blob is None:
            for rel in slide.part.rels.values():
                if "image" in rel.target_ref:
                    size = len(rel.target_part.blob)
                    if size > max_size:
                        max_size = size
                        largest_img_blob = rel.target_part.blob

        if largest_img_blob is None:
            print(f"Slide {slide_num}: No image found for OCR. Extracting raw text as fallback.")
            extracted_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    extracted_text.append(shape.text.strip())
            raw_text = "\n".join(extracted_text)
            with open(out_txt, "a", encoding="utf-8") as f:
                f.write(f"--- Slide {slide_num} ---\n{raw_text}\n\n")
            continue
            
        print(f"Slide {slide_num}: Sending image to Google Cloud Vision API...")
        try:
            ocr_text = detect_text(largest_img_blob)
            
            with open(out_txt, "a", encoding="utf-8") as f:
                f.write(f"--- Slide {slide_num} (Vision OCR) ---\n")
                f.write(f"{ocr_text.strip()}\n\n")

            print(f"Slide {slide_num}: ✅ Vision OCR completed.")
            
        except Exception as e:
            print(f"Slide {slide_num}: Vision API FAILED: {e}")

    print(f"\nAll slides processed. Final text saved to: {out_txt}")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python ocr_pptx.py <file.pptx>")
        sys.exit(1)
    process_pptx(sys.argv[1])
